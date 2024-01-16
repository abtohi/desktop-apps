import tkinter as tk
from ttkbootstrap.dialogs import Messagebox
import re
import pandas as pd
import os
from pptx import Presentation
import subprocess

class PPTExcel:
    def __init__(self, root, dtype, projectpath):
        self.root = root
        self.dtype = dtype
        self.projectpath = f'projects\\{projectpath.get()}'
    
    def extract_to_excel(self, **kwargs):
        selected_ppt = kwargs['ppt']
        progress_var = kwargs['prog_var']
        progress = kwargs['prog']
        lbl_status = kwargs['lbl_status']
        frame = kwargs['frame']
        progress_var.set(1)
        
        global extracted_text_dfs, extracted_table_dfs, extracted_chart_dfs

        ppt = Presentation(selected_ppt)
        progress_var.set(3)
        
        out = "output_data.xlsx"
        all_list = self.extract_all(ppt)
        self.save_to_excel(all_list, progress_var, out)
        lbl_status.config(text="All Data extracted and save to "+out)
        
        result = Messagebox.yesno(title="Options",message="Your data has been successfully extracted to "+out+"\n\nDo you want to open the output?",parent=frame)
        if result == 'Yes':
            output = f'{self.projectpath}\\output\\{out}'
            subprocess.run(['start', 'excel', output], shell=True)


        progress.pack_forget()
        lbl_status.config(text='')

    def identify_shape(self, shape):
        # Fungsi untuk mengidentifikasi tipe bentuk (shape)
        if shape.has_chart:
            return "Chart"
        elif shape.has_table:
            return "Table"
        elif shape.has_text_frame:
            return "Text"
        else:
            return "Unknown"
        
    def extract_texts(self, ppt):
        text_dfs = {}  # Dictionary untuk menyimpan DataFrame untuk setiap teks di slide yang berbeda
        for slide_number, slide in enumerate(ppt.slides, start=1):
            for shape in slide.shapes:
                text_data = []
                data_type = self.identify_shape(shape)
                if data_type == "Text":
                    text_frame = shape.text_frame
                    shape_name = shape.name

                    # Memisahkan teks berdasarkan baris baru (new line) jika ada
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text.strip()
                        text = text.replace("\x0b", "\n")
                        if text:
                            text_data.append([text])

                    # Buat DataFrame dari data teks
                    if text_data:
                        text_df = pd.DataFrame(text_data, columns=["Text"])
                        shape = f'{text_df.shape[0]}x{text_df.shape[1]}'
                        if text_df.shape[0] != 0:
                            text_dfs[f"Slide {slide_number}#{shape_name}#{shape}"] = text_df

        return text_dfs

    def extract_tables(self, ppt):
        table_number = 0
        table_dfs = {}  # Dictionary untuk menyimpan DataFrame untuk setiap table yang berbeda
        for slide_number, slide in enumerate(ppt.slides, start=1):
            for shape in slide.shapes:
                data_type = self.identify_shape(shape)
                if data_type == "Table":
                    table_number += 1  # Tambahkan nomor urutan untuk table yang ditemukan
                    shape_name = shape.name
                    table = shape.table
                    table_data = []
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(str(cell.text).replace('\x0b','\n'))
                        table_data.append(row_data)

                    # Buat DataFrame dari data table
                    columns = table_data[0]
                    table_df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    shape = f'{table_df.shape[0]}x{table_df.shape[1]}'

                    if table_df.shape[0] != 0:

                        for i, row in table_df.iterrows():
                            if all(cell == '' for cell in row):
                                if i == len(table_df)-1:
                                    table_df = table_df.drop(i)

                        count_true = sum(all(cell == '' for cell in row) for _, row in table_df.iterrows())
                        length = len(table_df)

                        if count_true != length:
                            # Tambahkan kolom Slide Number, Data Type, dan Table Number di kolom paling depan
                            table_dfs[f"Slide {slide_number}#{shape_name}#{shape}"] = table_df

        return table_dfs
    
    def extract_charts(self, ppt):
        chart_number = 0
        chart_dfs = {}  # Dictionary untuk menyimpan DataFrame untuk setiap chart yang berbeda
        for slide_number, slide in enumerate(ppt.slides, start=1):
            for shape in slide.shapes:
                data_type = self.identify_shape(shape)
                if data_type == "Chart":
                    chart = shape.chart
                    chart_number += 1  # Tambahkan nomor urutan untuk chart yang ditemukan

                    shape_name = shape.name
                    
                    data = []
                    for plot in chart.plots:
                        cat = []
                        for category in plot.categories.flattened_labels:
                            category = cat.append(category[0])
                        if len(cat) > 0:
                            for iseries, series in enumerate(plot.series, start=1):
                                for i, value in enumerate(series.values):
                                    if i != len(cat):
                                        index = i % len(cat)
                                        series_name = series.name or f'Column {iseries}'
                                        series_name = f"Column {iseries}" if series_name == "#REF!" else series_name
                                        new_cat = "Category" if len(cat[index]) == 0 else cat[index]
                                        new_values = round(value,2) if type(value) == float else value
                                        data.append([series_name, new_cat, new_values])

                    if chart.chart_type == 72:
                        data = []
                        next_idx = 1
                        for series in chart.series:
                            ser = series._ser
                            cat = series.name
                            x_pts = ser.xpath(".//c:xVal//c:pt")
                            for pt in x_pts:
                                str_value = pt.xpath("./c:v")[0].text
                                value = float(str_value)
                                idx = int(pt.get("idx"))+1
                                while next_idx < idx:
                                    data.append([cat, next_idx, None])
                                    next_idx += 1
                                data.append([cat, idx, value])
                                next_idx += 1

                    # Buat DataFrame dari data yang telah di-stack
                    chart_df = pd.DataFrame(data, columns=["Legend", "Categories", "Values"])

                    # Mengganti nama indeks dan kolom
                    shape = f'{chart_df.shape[0]}x{chart_df.shape[1]}'
                    if chart_df.shape[0] != 0:
                        chart_dfs[f"Slide {slide_number}#{shape_name}#{shape}#{chart.chart_type}"] = chart_df

        return chart_dfs
    
    def extract_all(self, ppt):
        all_data = []
        if 'Texts' in self.dtype:
            texts = self.extract_texts(ppt)
            all_data.append(texts)
        if 'Tables' in self.dtype:
            tables = self.extract_tables(ppt)
            all_data.append(tables)
        if 'Charts' in self.dtype:
            charts = self.extract_charts(ppt)
            all_data.append(charts)

        return all_data

    def save_to_excel(self, dataframes, progress_var, file_name):
        path = self.projectpath
        path = f'{str(path)}\output'
        if not os.path.exists(path):
            os.makedirs(path)
        
        with pd.ExcelWriter(str(path)+'\\'+file_name, engine='xlsxwriter') as writer:
            for index, (dfs, datatype) in enumerate(zip(dataframes, self.dtype)):
                row_usage = 1

                progress_var.set(5)
                for i, (table_name, df) in enumerate(dfs.items(),start=1):
                    progress_var.set((index * (100/len(dataframes))) + ((i / len(dfs.items())) * (100/len(dataframes))))
                    self.root.update_idletasks()
                    df.to_excel(writer, sheet_name=datatype, index=False, startrow=row_usage-1, startcol=1, header=True)
                
                    worksheet = writer.sheets[datatype]
                    worksheet.write_string(row_usage-1, 0, table_name)  # Menulis nama lembar ke kolom pertama
                    row_usage += len(df) + 2
                
    def index_to_excel_column(self, index):
        result = ""
        while index > 0:
            remainder = (index - 1) % 26
            result = chr(ord('A') + remainder) + result
            index = (index - 1) // 26
        return result
    
    def clean_data(self, input_string):
        # Menghapus spasi dan simbol-simbol non-alfanumerik
        cleaned_string = re.sub(r'[^a-zA-Z0-9]', '', input_string)
        # Mengubah hasilnya menjadi huruf kecil
        cleaned_string = cleaned_string.lower() 
        return cleaned_string