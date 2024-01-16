import pandas as pd
import re
import os
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from ttkbootstrap.dialogs import Messagebox
import datetime
import subprocess

class ExcelPPT:
    def __init__(self, dtype, lbl_status, root, selected_ppt, projectpath):
        self.dtype = dtype
        self.lbl_status = lbl_status
        self.root = root
        self.logreport = []
        self.pptx_file = selected_ppt or ""
        self.projectpath = f'projects/{projectpath.get()}/output'
        self.prs = Presentation(self.pptx_file)

    # Function to replace_ppt_values from the selected Excel file
    def replace_ppt_values(self, **kwargs):
        progress = kwargs['prog'] 
        progress_var = kwargs['prog_var']
        frame = kwargs['frame']

        output_filename = 'output_data.xlsx'
        selected_excel = f"{self.projectpath}/{output_filename}"
        self.modify_all(selected_excel, progress_var)
        report = pd.DataFrame(self.logreport)
        report.to_csv(f'{self.projectpath}/log.csv')

        # Save the modified presentation with a new filename
        new_pptx_file = f'{self.projectpath}\\{os.path.splitext(os.path.basename(self.pptx_file))[0]}_modified.pptx'
        self.prs.save(new_pptx_file)
        self.lbl_status.config(text="Your powerpoint has been successfully modified")
        result = Messagebox.yesno(title="Options",message="Your powerpoint has been successfully modified to "+new_pptx_file+"\n\nDo you want to open the output?", parent=frame)
        if result == 'Yes':
            os.system(f'start powerpnt "{new_pptx_file}"')
        else:
            print("Your powerpoint has been successfully modified to "+new_pptx_file)
        
        progress.pack_forget()
        
    def modify_texts(self, selected_excel, progress_var):
        self.lbl_status.config(text="Modifying Texts")
        wb = load_workbook(selected_excel)
        data = []
        for item in wb:
            if item.title == 'Texts':
                for index, value in enumerate(item.values, start=1):
                    progress_var.set((index / len(list(item.values))) * 100/2)
                    self.root.update_idletasks()
                    nrow = 0
                    ncol = 0
                    if value[0] != None:
                        shape = str(value[0]).split('#')[2]
                        nrow = int(str(shape).split('x')[0])
                        ncol = int(str(shape).split('x')[1])

                    if nrow != 0 and ncol != 0:
                        df = pd.read_excel(selected_excel, sheet_name='Texts', nrows=nrow, skiprows=index-1, usecols=range(1,ncol+1))
                        val = {'Identity':value[0], 'Data':df}
                        data.append(val)
        
        if self.pptx_file != '':
            for index, slide in enumerate(self.prs.slides, start=1):
                progress_var.set((100/(len(self.dtype))) + ((index / len(self.prs.slides)) * 100/2))
                self.root.update_idletasks()

                for shape in slide.shapes:
                    shape_name = shape.name
                    
                    for entry in data:
                        idname = f'Slide {index}#{shape_name}'
                        identity = f"{str(entry['Identity']).split('#')[0]}#{str(entry['Identity']).split('#')[1]}"
                        if idname == identity:
                            df = entry['Data']
                            current_time = datetime.datetime.now()
                            timestamp = current_time.strftime("%Y-%m-%d %H:%M:%S")

                            #Modifikasi Text
                            if shape.has_text_frame:
                                for id_par, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                                    num_rows, num_cols = df.shape
                                    for i in range(num_rows):
                                        if i+1 == id_par:
                                            new_value = str(df.iloc[i,0])

                                            split_value = new_value.split('\n')
                                            try:
                                                if len(split_value) > 1:
                                                    for new_text, run in zip(split_value, paragraph.runs):
                                                        font_name = run.font.name
                                                        font_size = run.font.size
                                                        font_bold = run.font.bold
                                                        font_italic = run.font.italic
                                                        font_underline = run.font.underline

                                                        run.font.name = font_name
                                                        run.font.size = font_size
                                                        run.font.underline = font_underline
                                                        run.font.bold = font_bold
                                                        run.font.italic = font_italic
                                                        run.text = new_text
                                                else:
                                                    # Update via Paragraph
                                                    if len(paragraph.runs) > 0:
                                                        font_name = paragraph.runs[0].font.name 
                                                        font_size = paragraph.runs[0].font.size
                                                        font_bold = paragraph.runs[0].font.bold
                                                        font_italic = paragraph.runs[0].font.italic
                                                        font_underline = paragraph.runs[0].font.underline

                                                    paragraph.text = new_value

                                                    if len(paragraph.runs) > 0:
                                                        paragraph.runs[0].font.name = font_name
                                                        paragraph.runs[0].font.size = font_size
                                                        paragraph.runs[0].font.bold = font_bold
                                                        paragraph.runs[0].font.italic = font_italic
                                                        paragraph.runs[0].font.underline = font_underline

                                                log = {"slide_number":index,"shape":shape_name,"type":"","status":"Success", "message":"","timestamp":timestamp}
                                            
                                            except Exception as e:
                                                msg_error = str(e)
                                                log = {"slide_number":index,"shape":shape_name,"type":"","status":"Error", "message":msg_error,"timestamp":timestamp}
                                            
                                            self.logreport.append(log)

    def modify_tables(self, selected_excel, progress_var):
        self.lbl_status.config(text="Modifying Tables")
        wb = load_workbook(selected_excel)
        data = []
        for item in wb:
            if item.title == 'Tables':
                for index, value in enumerate(item.values, start=1):
                    progress_var.set(index / len(list(item.values)) * 100/2)
                    self.root.update_idletasks()
                    nrow = 0
                    ncol = 0
                    if value[0] != None:
                        shape = str(value[0]).split('#')[2]
                        nrow = int(str(shape).split('x')[0])
                        ncol = int(str(shape).split('x')[1])

                    if nrow != 0 and ncol != 0:
                        df = pd.read_excel(selected_excel, sheet_name='Tables', nrows=nrow, skiprows=index-1, usecols=range(1,ncol+1))
                        val = {'Identity':value[0], 'Data':df}
                        data.append(val)

        if self.pptx_file != '':
            for index, slide in enumerate(self.prs.slides, start=1):
                progress_var.set(100/(len(self.dtype)) + ((index / len(self.prs.slides)) * 100/2))
                self.root.update_idletasks()

                for shape in slide.shapes:
                    shape_name = shape.name
                    
                    for entry in data:
                        idname = f'Slide {index}#{shape_name}'
                        identity = f"{str(entry['Identity']).split('#')[0]}#{str(entry['Identity']).split('#')[1]}"
                        if idname == identity:
                            df = entry['Data']
                            current_time = datetime.datetime.now()
                            timestamp = current_time.strftime("%Y-%m-%d %H:%M:%S")

                            #Modifikasi Tabel
                            if shape.has_table:
                                table = shape.table
                                num_rows, num_cols = df.shape
                                try:
                                    if len(table.rows) - 1 == num_rows and len(table.columns) == num_cols:
                                        last_row = None
                                        last_col = None
                                        for i in range(num_rows):
                                            for j in range(num_cols):
                                                cell_value = str(df.iloc[i, j])
                                                if cell_value == 'nan': continue
                                                split_value = cell_value.split('\n')

                                                # Patch Style and Font
                                                paragraph = table.cell(i+1, j).text_frame.paragraphs
                                                for new_text, par in zip(split_value, paragraph):
                                                    for run in par.runs:
                                                        font_name = run.font.name
                                                        font_size = run.font.size
                                                        font_bold = run.font.bold
                                                        font_italic = run.font.italic
                                                        font_underline = run.font.underline

                                                        run.font.name = font_name
                                                        run.font.size = font_size
                                                        run.font.underline = font_underline
                                                        run.font.bold = font_bold
                                                        run.font.italic = font_italic
                                                        
                                                        if (last_row, last_col) == (i,j) and i != 0:
                                                            pass
                                                        else:
                                                            last_row = i
                                                            last_col = j
                                                            run.text = new_text
                                          
                                    log = {"slide_number":index,"shape":shape_name,"type":"","status":"Success", "message":"","timestamp":timestamp}
                                            
                                except Exception as e:
                                    msg_error = str(e)
                                    log = {"slide_number":index,"shape":shape_name,"type":"","status":"Error", "message":msg_error,"timestamp":timestamp}
                                
                                self.logreport.append(log)

    def modify_charts(self, selected_excel, progress_var):
        self.lbl_status.config(text="Modifying Charts")
        wb = load_workbook(selected_excel)
        data = []
        for item in wb:
            if item.title == 'Charts':
                for index, value in enumerate(item.values, start=1):
                    progress_var.set(index / len(list(item.values)) * 100/2)
                    self.root.update_idletasks()
                    nrow = 0
                    ncol = 0
                    if value[0] != None:
                        shape = str(value[0]).split('#')[2]
                        nrow = int(str(shape).split('x')[0])
                        ncol = int(str(shape).split('x')[1])

                    if nrow != 0 and ncol != 0:
                        df = pd.read_excel(selected_excel, sheet_name='Charts', nrows=nrow, skiprows=index-1, usecols=range(1,ncol+1))
                        val = {'Identity':value[0], 'Data':df}
                        data.append(val)
        
        if self.pptx_file != '':
            for index, slide in enumerate(self.prs.slides, start=1):
                progress_var.set(100/(len(self.dtype)) + ((index / len(self.prs.slides)) * 100/2))
                self.root.update_idletasks()

                for shape in slide.shapes:
                    shape_name = shape.name
                    
                    for entry in data:
                        idname = f'Slide {index}#{shape_name}'
                        identity = f"{str(entry['Identity']).split('#')[0]}#{str(entry['Identity']).split('#')[1]}"
                        if idname == identity:
                            df = entry['Data']
                            current_time = datetime.datetime.now()
                            timestamp = current_time.strftime("%Y-%m-%d %H:%M:%S")

                            #Modifikasi Chart
                            if shape.has_chart:
                                chart = shape.chart
                                
                                # Menentukan data yang akan digunakan dalam chart
                                if str(idname).split('#')[0].replace('Slide ','') == str(index):
                                    jenis = str(entry['Identity']).split('#')[3]
                                    chart_type = jenis.replace('(','').replace(')','')
                                    chart_type = chart_type.split(' ')[1]
                                    chart_type = int(chart_type)
                                    series = tuple(df.iloc[:,2])

                                    if len(series) == 0:
                                        continue

                                    categories = tuple(df.iloc[:,1])        
                                    chart_data = CategoryChartData()

                                    report = self.charttype_recognizing(
                                        chart_type=chart_type, 
                                        chart_data=chart_data, 
                                        categories=categories, 
                                        df=df, chart=chart)
                                    
                                    status = report['status']
                                    message = report['message']
                                    log = {"slide_number":index,"shape":shape_name,"type":jenis,"status":status, "message":message, "timestamp":timestamp}
                                    self.logreport.append(log)

    def modify_all(self, selected_excel, progress_var):
        if 'Texts' in self.dtype:
            self.modify_texts(selected_excel, progress_var)
        if 'Tables' in self.dtype:
            self.modify_tables(selected_excel, progress_var)
        if 'Charts' in self.dtype:
            self.modify_charts(selected_excel, progress_var)

    def convert_range(self, excel_range):
        """ Converts a range (i.e. something like 'A3:D20') and returns
        the corresponding arguments to use in pd.read_excel."""

        # Get cell addresses from range (i.e. A3 and D20)
        upper_left, lower_right = excel_range.split(':')

        # Convert cell address ('A3') to col (A) and row (3)
        left_col, top_row = list(filter(None, re.split('(\d+)', upper_left)))
        right_col, bottom_row = list(filter(None, re.split('(\d+)', lower_right )))

        return {'usecols': f'{left_col}:{right_col}',
                'skiprows': int(top_row) - 1,
                'nrows': int(bottom_row) - int(top_row) + 1}
    
    def clean_data(self, input_string):
        # Menghapus spasi dan simbol-simbol non-alfanumerik
        cleaned_string = re.sub(r'[^a-zA-Z0-9]', '', input_string)
        # Mengubah hasilnya menjadi huruf kecil
        cleaned_string = cleaned_string.lower() 
        return cleaned_string
    
    def open_log(self):
        out = f'{self.projectpath}\log.csv'
        subprocess.run(['start', 'excel', out], shell=True)
    

    def charttype_recognizing(self, **kwargs):
        df = kwargs['df']
        categories = kwargs['categories']
        chart_type = kwargs['chart_type']
        chart_data = kwargs['chart_data']
        chart = kwargs['chart']
        
        try:
            if chart_type == 4 or chart_type == 5 or chart_type == 51 or chart_type == -4120 or chart_type == 72: #Column Clustered, doughnut
                chart_data.categories = categories
                val = df.iloc[:,2].fillna('')
                chart_data.add_series("Series", val)
                chart.replace_data(chart_data)

            elif chart_type == 52 or chart_type == 53 or chart_type == 57 or chart_type == 58: #Bar Clustered
                chart_data.categories = tuple(df.iloc[:,1].unique())
                for leg in df['Legend'].unique():
                    val = df[df['Legend'] == leg]
                    val = tuple(val.iloc[:,2].fillna(''))
                    chart_data.add_series(leg, val)
                chart.replace_data(chart_data)
            
            else:
                status = 'Skipped'
                message = f'Chart type {chart_type} is not defined yet'
            
            try:
                status = status
                message = message
            except:
                status = 'Success'
                message = ''
            
            return {'type':chart_type, 'status':status, 'message': message}
        
        except Exception as e:
            error_message = str(e)
            return {'type':chart_type, 'status':'Error', 'message':error_message}
            
        
