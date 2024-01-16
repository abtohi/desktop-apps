import os
import tkinter as tk
from tkinter import filedialog
import ttkbootstrap as tb
from ttkbootstrap import Style, font
from ttkbootstrap.constants import *
from ttkbootstrap.dialogs import Messagebox
from PIL import ImageTk, Image
from pptx import Presentation
from pandastable import Table, config
import pandas as pd

from excel_ppt import ExcelPPT
from ppt_excel import PPTExcel
from convert_data import ConvertData
from sync import Synchronize
from reconvert_data import ReConvert
from display.menu import Menu

#Kelas display, untuk tampilan aplikasi
class Display:
    def __init__(self):
        self.extracted_text_dfs = {}
        self.extracted_table_dfs = {}
        self.extracted_chart_dfs = {}

    def select_ppt(self):
        if var_new_project.get() == '':
            Messagebox.show_warning(message="Please create or select your project", parent=frame4)
        else:
            file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
            text_path.set(file_path)
            self.selected_ppt = file_path
            prs = Presentation(file_path)

            num_slide = 0
            num_shape = 0
            num_text = 0
            num_table = 0
            num_chart = 0

            for slide in prs.slides:
                num_slide += 1
                for shape in slide.shapes:
                    num_shape += 1
                    if shape.has_text_frame:
                        text = shape.text
                        if len(text) !=0:
                            num_text += 1
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text)
                            table_data.append(row_data)
                        table_df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        if len(table_df):
                            num_table +=1
                    if shape.has_chart:
                        num_chart += 1
            
            detail_sum = ("Slide","Shape","Text","Table","Chart")
            value_sum = (num_slide, num_shape, num_text, num_table, num_chart)

            for index, (det,val) in enumerate(zip(detail_sum, value_sum)):
                if index == 0:
                    lbl_sum1 = tb.Label(frame6,text=f"Total {det}", bootstyle="primary")
                    lbl_sum1.grid(row=index, column=0, sticky="w", padx=18, pady=(19,0))

                    lbl_sum2 = tb.Label(frame6,text=":", bootstyle="primary")
                    lbl_sum2.grid(row=index, column=1, pady=(19,0))

                    lbl_sum3 = tb.Label(frame6,text=val, bootstyle="primary")
                    lbl_sum3.grid(row=index, column=2, sticky="w", padx=4, pady=(19,0))
                else:
                    lbl_sum1 = tb.Label(frame6,text=f"Total {det}", bootstyle="primary")
                    lbl_sum1.grid(row=index, column=0, sticky="w", padx=18)

                    lbl_sum2 = tb.Label(frame6,text=":", bootstyle="primary")
                    lbl_sum2.grid(row=index, column=1)

                    lbl_sum3 = tb.Label(frame6,text=val, bootstyle="primary")
                    lbl_sum3.grid(row=index, column=2, sticky="w", padx=4)

    def select_data_table(self):
        file_path2 = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])
        text_path2.set(file_path2)

    def process(self):
        dtype_selected = result_type.get().split('|')
        if text_path.get():
            if actions_var.get():
                if dtype_selected != ['']:
                    if actions_var.get() == "Extract":
                        pptexcel = PPTExcel(root, dtype_selected, var_new_project)
                        progressbar.pack(ipady=scaled(5))
                        pptexcel.extract_to_excel(
                            ppt = self.selected_ppt, 
                            prog_var = progress_var, 
                            prog = progressbar, 
                            lbl_status = lbl_status,
                            frame = frame4)
                        
                    elif actions_var.get() == "Convert":
                        progressbar.pack(ipady=scaled(5))
                        convertdata = ConvertData(
                            projectpath = var_new_project,  
                            lbl_status = lbl_status, 
                            dtype=dtype_selected, 
                            root = root,
                            frame4 = frame4)
                            
                        convertdata.convert(progress_var, progressbar)
                    
                    elif actions_var.get() == "Synchronize":
                        selected = country.get()

                        if selected == "Indonesia":
                            progressbar.pack(ipady=scaled(5))
                            progress_var.set(1)
                            root.update_idletasks()
                            if text_path2.get():
                                #Sync
                                synchronize = Synchronize(root, lbl_status, dtype_selected, progressbar, progress_var, frame4, text_path2.get(), var_new_project)
                                synchronize.sync()
                            else:
                                Messagebox.show_warning(title="Warning",message="Please select your data table first")
                        else:
                            Messagebox.show_warning(title="Warning",message="The country you selected is currently unavailable")
                    
                    elif actions_var.get() == "Reconvert":
                        progressbar.pack(ipady=scaled(5))
                        progress_var.set(1)
                        root.update_idletasks()
                        re_convert = ReConvert(root, dtype_selected, progressbar, progress_var, frame4, lbl_status, var_new_project)
                        re_convert.reconvert()

                    elif actions_var.get() == "Modify PPT": 
                        path = f'projects/{var_new_project.get()}/output'
                        lists_dir = os.listdir(path)
                        status = 1
                        output_filename = 'output_data.xlsx'

                        if output_filename not in lists_dir :
                            Messagebox.show_error(title="Error",message="There is no output file that you extracted", parent=frame4)
                            status = 0

                        if status == 1:
                            progressbar.pack(ipady=5)
                            excelppt = ExcelPPT(dtype_selected, lbl_status, root, self.selected_ppt, var_new_project)
                            excelppt.replace_ppt_values(prog = progressbar, prog_var = progress_var, frame = frame4)
                            lbl_log.config(text="Open Log")
                            lbl_log.bind("<Button-1>", lambda event: excelppt.open_log())
                else:
                    Messagebox.show_warning(title="Warning",message="Please select your data type", parent=frame4)
            else:
                Messagebox.show_warning(title="Warning",message="Please select your action", parent=frame4)
        else:
            Messagebox.show_warning(title="Warning", message="Please select your PPT File", parent=frame4)

    def view_data(self):
        try:
            if text_path.get():
                dtype_selected = result_type.get().split('|')
                if dtype_selected != ['']:
                    app = tb.Toplevel(title="Data View")
                    app.geometry(f"{scaled(1010)}x{scaled(500)}")

                    titlenb = tb.Label(app,text="Select Data Type")
                    titlenb.place(x=scaled(20), y=scaled(20))
                    option = ["Texts","Tables","Charts"]
                    
                    combobox = tb.Combobox(app, values=option, bootstyle="primary", width=scaled(10), height=scaled(5))
                    combobox.current(option.index(dtype_selected[0]))
                    combobox.place(x=scaled(130), y=scaled(17))

                    def option_selected(event):
                        option_selected = combobox.get()
                        if option_selected == "Texts":
                            self.extracted_text_dfs = pptexcel.extract_texts(ppt)
                            self.display_structured_tables(self.extracted_text_dfs, nb)
                        elif option_selected == "Tables":
                            self.extracted_table_dfs = pptexcel.extract_tables(ppt)
                            self.display_structured_tables(self.extracted_table_dfs, nb)
                        elif option_selected == "Charts":
                            self.extracted_chart_dfs = pptexcel.extract_charts(ppt)
                            self.display_structured_tables(self.extracted_chart_dfs, nb)
                    
                    combobox.bind("<<ComboboxSelected>>", option_selected)

                    nb = tb.Notebook(app, bootstyle="info")
                    nb.place(x=scaled(20), y=scaled(65))
                    
                    ppt = Presentation(self.selected_ppt)
                    pptexcel = PPTExcel(root, dtype_selected, var_new_project)
                    if "Texts" in dtype_selected:
                        self.extracted_text_dfs = pptexcel.extract_texts(ppt)
                        self.display_structured_tables(self.extracted_text_dfs, nb)
                    elif "Tables" in dtype_selected:
                        self.extracted_table_dfs = pptexcel.extract_tables(ppt)
                        self.display_structured_tables(self.extracted_table_dfs, nb)
                    elif "Charts" in dtype_selected:
                        self.extracted_chart_dfs = pptexcel.extract_charts(ppt)
                        self.display_structured_tables(self.extracted_chart_dfs, nb)
                    else:
                        pass
                    app.mainloop()
        
                else:
                    Messagebox.show_warning(title="Warning",message="Please select your data type", parent=frame4)
            else:
                Messagebox.show_warning(title="Warning", message="Please select your PPT File", parent=frame4)
        except Exception as e:
            msg_error = str(e)
            Messagebox.show_error(title="Error",message=msg_error, parent=frame4)

    def display_structured_tables(self, dataframes, notebook):
        for tab in notebook.tabs():
            notebook.forget(tab)

        for df_name, df in dataframes.items():
            df_frame = tb.Frame(notebook)
            table = Table(df_frame, dataframe=df, showtoolbar=True, showstatusbar=True, enable_menus=False)
            table.show()
            notebook.add(df_frame, text=df_name)
            
            options=config.load_options()
            options={'fontsize':10,'cellwidth': 80, 'grid_color':'#828485', 'colheadercolor':'#438cb0'}
            
            config.apply_options(options,table)

        if dataframes:
            notebook.select(0)

    def activate_data_table(self):
        if actions_var.get() == "Synchronize":
            browse_btn2.config(state=tb.NORMAL)
            frame45.config(bootstyle="primary")
            dt_path.configure(state=tb.NORMAL)
        else:
            browse_btn2.config(state=tb.DISABLED)
            frame45.config(bootstyle="secondary")
            dt_path.configure(state=tb.DISABLED)

    def create_folder(self,foldername):
        path = f'{os.getcwd()}\projects\{var_new_project.get()}'
        path = f'{str(path)}\{foldername}'
        if not os.path.exists(path):
            os.makedirs(path)
        return path

    def open_folder(self):
        if var_new_project.get() == '':
            Messagebox.show_warning(message="Please create or select your project", parent=frame4)
        else:
            path = self.create_folder('output')
            os.startfile(path)
        
    def close_window(self):
        root.destroy()
    
    def dtype_selected(self):
        result_type.set('|'.join([dtype for dtype, var in zip(data_type, data_type_var) if var.get()]))

    def select_country(self):
        if actions_var.get() == "Synchronize":
            npframe = tb.Frame(root, border=1, relief="solid")
            npframe.place(x=scaled(40), y=scaled(100), width=530, height=200)

            labeltitle = tb.Label(npframe, text="Select Country", font=("Arial Bold",15), bootstyle="primary")
            labeltitle.place(x=scaled(185), y=scaled(20))

            labelinfo = tb.Label(npframe, text="Please select a country to determine the format of the table data", font=("Helvetica",10,"italic"), bootstyle="primary")
            labelinfo.place(x=scaled(75), y=scaled(50))

            labelproject = tb.Label(npframe, text="Country",font=("Helvetica",11), bootstyle="primary")
            labelproject.place(x=scaled(125), y=scaled(90))

            option = ["Indonesia","Malaysia","Thailand"]
            combobox = tb.Combobox(npframe, values=option, bootstyle="primary")
            combobox.place(x=scaled(190), y=scaled(87), width=scaled(200), height=scaled(25))
            combobox.current(option.index(option[0]))

            selectbtn = tb.Button(npframe, text="Select", command=lambda: (country.set(combobox.get()), npframe.destroy()))
            selectbtn.place(x=scaled(170), y=scaled(142), width=90)

            cancelbtn = tb.Button(npframe, text="Cancel", bootstyle="danger", command=lambda: (npframe.destroy()))
            cancelbtn.place(x=scaled(280), y=scaled(142), width=90)

if __name__ == "__main__":
    app = Display()
    root = tb.Window()

    style = Style(theme="litera")

    root.title("Powerpoint Extractor v1.5")
    root.resizable(width=False, height=True)

    dpi_ori  = 96
    dpi_user = root.winfo_fpixels('1i')
    scaled = lambda x: round(x * (dpi_user / dpi_ori))
    
    root.geometry(f"{scaled(610)}x{scaled(750)}")
    root.place_window_center()
    width = root.winfo_width()
    height = root.winfo_height()
    selected_ppt = None

    width_default = scaled(572)
    height_default = scaled(112)

    frame1 = tb.Frame(root)
    frame1.place(x=scaled(19), y=scaled(10), width=width_default, height=height_default)

    ico_path = os.path.join(os.path.dirname(__file__), "img", "ext.ico")
    root.iconbitmap(ico_path)

    img_path1 = os.path.join(os.path.dirname(__file__), "img", "NIQ Logo.jpg")
    img_path2 = os.path.join(os.path.dirname(__file__), "img", "NIQ Vision Logo.jpg")

    img1 = ImageTk.PhotoImage(Image.open(img_path1).resize((height_default, height_default)))
    img2 = ImageTk.PhotoImage(Image.open(img_path2).resize((scaled(450),height_default)))

    logo = tb.Label(frame1, image=img1)
    cover = tb.Label(frame1, image=img2)

    logo.grid(row=0, column=0)
    cover.grid(row=0, column=1)

    frame2 = tb.Frame(root)
    frame2.place(x=scaled(19), y=scaled(130), width=width_default, height=height_default)

    title = tb.Label(frame2, text="Power Point Extractor", bootstyle="primary")
    title.config(font=("Arial Bold", 16))
    title.pack(pady=scaled(15))

    labelfont = font.Font(family='TkDefaultFont',size=9, slant="italic")
    labelproject = tb.Label(root, text=f"Project Name : ", bootstyle="secondary", font=labelfont)
    labelproject.place(x=scaled(19), y=scaled(183))

    frame3 = tb.LabelFrame(root,text="Select Your Powerpoint", bootstyle="primary")
    frame3.place(x=scaled(19), y=scaled(210), width=width_default, height=scaled(80))

    text_path = tb.StringVar()
    ppt_path = tb.Entry(root, bootstyle="default", textvariable=text_path)
    ppt_path.configure(state="readonly")
    ppt_path.place(x=scaled(45), y=scaled(240), width=scaled(420))

    browse_btn = tb.Button(root, text="Browse",bootstyle="info", command=app.select_ppt)
    browse_btn.place(x=scaled(480), y=scaled(240), width=scaled(85))

    frame4 = tb.LabelFrame(root,text="Select Your Action", bootstyle="primary")
    frame4.place(x=scaled(19), y=scaled(305), width=width_default, height=scaled(74))

    var_new_project = tb.StringVar()
    mainmenu = Menu(app,root, var_new_project, frame4, style, scaled, labelproject)
    mainmenu.menu()

    actions_var = tb.StringVar()
    actions = ('Extract', 'Convert', 'Synchronize', 'Reconvert', 'Modify PPT')

    country = tb.StringVar()

    grid_frame4 = 0
    for index, action in enumerate(actions):
        radio = tb.Radiobutton(frame4, text=action, value=action, variable=actions_var, command=lambda: (lbl_log.config(text=""), app.activate_data_table(), app.select_country()))
        if grid_frame4 == 0:
            radio.grid(column=grid_frame4, row=0, padx=(scaled(65),scaled(10)), pady=scaled(20))
        elif grid_frame4 == 1:
            radio.grid(column=grid_frame4, row=0, padx=scaled(10), pady=scaled(10))
        elif grid_frame4 ==2:
            radio.grid(column=grid_frame4, row=0, padx=scaled(10), pady=scaled(10))
        elif grid_frame4 ==3:
            radio.grid(column=grid_frame4, row=0, padx=scaled(10), pady=scaled(10))
        else:
            radio.grid(column=grid_frame4, row=0, padx=(scaled(10),scaled(65)), pady=scaled(20))
        grid_frame4 += 1

    frame45 = tb.LabelFrame(root, text="Select Data Table", bootstyle="secondary")
    frame45.place(x=scaled(19), y=scaled(395), width=scaled(350), height=scaled(80))

    text_path2 = tb.StringVar()
    dt_path = tb.Entry(root, bootstyle="default", textvariable=text_path2)
    dt_path.configure(state=tb.DISABLED)
    dt_path.place(x=scaled(45), y=scaled(425), width=scaled(200))

    browse_btn2 = tb.Button(root, text="Browse",bootstyle="info", command=app.select_data_table, state=tb.DISABLED)
    browse_btn2.place(x=scaled(260), y=scaled(425), width=scaled(85))

    frame5 = tb.LabelFrame(root, text="Select Data Type", bootstyle="primary")
    frame5.place(x=scaled(19), y=scaled(485), width=scaled(350), height=scaled(60))
    
    data_type = ("Texts","Tables","Charts")
    result_type = tb.StringVar()

    data_type_var = [tk.BooleanVar() for _ in data_type]

    grid_frame5 = 0
    for dtype, var in zip(data_type, data_type_var):
        check = tb.Checkbutton(frame5, text=dtype, variable=var, bootstyle="primary-round-toggle", command=lambda: (app.dtype_selected(),lbl_log.config(text=""), lbl_status.config(text="")))
        if grid_frame5 == 0:
            check.grid(column=grid_frame5, row=0, padx=(scaled(70), scaled(10)), pady=scaled(10), sticky="w")
        elif grid_frame5 == 1:
            check.grid(column=grid_frame5, row=0, padx=scaled(10), pady=scaled(10), sticky="w")
        else:
            check.grid(column=grid_frame5, row=0, padx=scaled(10), pady=scaled(10), sticky="w")
        grid_frame5 += 1

    frame6 = tb.LabelFrame(root, text="Summary", bootstyle="primary")
    frame6.place(x=scaled(385), y=scaled(395), width=scaled(206), height=scaled(150))

    frame7 =  tb.LabelFrame(root,text="Execute Your Action", bootstyle="primary")
    frame7.place(x=scaled(19), y=scaled(564), width=width_default, height=scaled(80))

    btn_process = tb.Button(root, text="Process",bootstyle="primary", command=app.process)
    btn_process.place(x=scaled(75), y=scaled(595), width=scaled(100))

    btn_view = tb.Button(root, text="View Data",bootstyle="success", command=app.view_data)
    btn_view.place(x=scaled(195), y=scaled(595), width=scaled(100))

    btn_folder = tb.Button(root, text="Open Folder",bootstyle="success", command=app.open_folder)
    btn_folder.place(x=scaled(315), y=scaled(595), width=scaled(100))

    btn_cancel = tb.Button(root, text="Cancel",bootstyle="danger", command=app.close_window)
    btn_cancel.place(x=scaled(435), y=scaled(595), width=scaled(100))

    frame8 =  tb.Frame(root)
    frame8.place(x=scaled(19), y=scaled(655), width=width_default, height=scaled(80))

    frame9 = tb.Frame(root)
    frame9.place(x=scaled(19), y=scaled(655), width=width_default, height=scaled(70))

    lbl_status = tb.Label(frame9, text="")
    lbl_status.pack()

    lbl_log = tb.Label(frame9, text="", cursor="hand2", font=('Helvetica',9,'bold','underline'), foreground="blue")
    lbl_log.pack()

    # default striped progressbar style
    progress_var = tk.DoubleVar()
    progressbar = tb.Progressbar(frame9, bootstyle="success", length=width_default, mode="determinate", variable=progress_var)

    root.mainloop()